# -*- coding: utf-8 -*-
"""
Create a presentation deck for Problem 4 speech signal processing.

The deck uses a warm off-white Morandi palette and creates PPT-friendly
derivative charts from the experiment outputs.
"""

from __future__ import annotations

import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import soundfile as sf
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from scipy import signal


ROOT = Path(__file__).resolve().parents[1]
AUDIO = ROOT / "audio"
IMAGES = ROOT / "images"
ASSETS = ROOT / "ppt_assets"
OUT = ROOT / "第四问_语音信号处理汇报.pptx"

W, H = 13.333, 7.5

BG = "F7F1E5"
PANEL = "FFFCF5"
TEXT = "3F3A36"
MUTED = "7B746C"
BLUE = "9DB7C7"
BLUE_DARK = "6F8FA3"
SAGE = "A8B8A0"
PINK = "C9A7A2"
TERR = "C8A27A"
WARM_GRAY = "B9B2A7"
LINE = "DED5C8"


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def mpl_color(hex_color: str) -> str:
    return f"#{hex_color.strip('#')}"


def ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)


def read_audio(path: Path) -> tuple[np.ndarray, int]:
    data, sr = sf.read(path, always_2d=True)
    return data.astype(np.float32).mean(axis=1), int(sr)


def save_fig(path: Path) -> None:
    plt.savefig(path, bbox_inches="tight", facecolor=mpl_color(PANEL), dpi=220)
    plt.close()


def morandi_style() -> None:
    plt.rcParams.update(
        {
            "font.sans-serif": ["Microsoft YaHei", "SimHei", "Calibri"],
            "axes.unicode_minus": False,
            "figure.facecolor": mpl_color(PANEL),
            "axes.facecolor": mpl_color(PANEL),
            "axes.edgecolor": "#CFC6BA",
            "axes.labelcolor": mpl_color(TEXT),
            "xtick.color": mpl_color(MUTED),
            "ytick.color": mpl_color(MUTED),
            "text.color": mpl_color(TEXT),
            "grid.color": "#E7DDD0",
        }
    )


def create_metric_charts() -> dict[str, Path]:
    morandi_style()
    metrics = pd.read_csv(IMAGES / "problem4_speech_metrics.csv")
    order = ["Butterworth", "Chebyshev-I", "Elliptic", "FIR-Hamming", "ReferenceGuidedFilter", "DifferentiableFilter", "TrainableRefGuided"]
    metrics["method"] = pd.Categorical(metrics["method"], categories=order, ordered=True)
    metrics = metrics.sort_values(["method", "stem"])

    colors = {
        "Butterworth": "#B9B2A7",
        "Chebyshev-I": "#C8A27A",
        "Elliptic": "#C9A7A2",
        "FIR-Hamming": "#A8B8A0",
        "ReferenceGuidedFilter": "#9DB7C7",
        "DifferentiableFilter": "#6F8FA3",
        "TrainableRefGuided": "#7F9FAD",
    }

    pivot_sisdr = metrics.pivot(index="method", columns="stem", values="si_sdr_db").loc[order]
    fig, ax = plt.subplots(figsize=(8.6, 4.6))
    x = np.arange(len(order))
    ax.bar(x - 0.18, pivot_sisdr["male"], width=0.36, color=[colors[m] for m in order], alpha=0.95, label="male")
    ax.bar(x + 0.18, pivot_sisdr["female"], width=0.36, color=[colors[m] for m in order], alpha=0.55, label="female")
    ax.axhline(0, color="#8E857A", linewidth=1)
    ax.set_xticks(x)
    ax.set_xticklabels(["Butter", "Cheby-I", "Ellip", "FIR", "Ref", "Diff", "T-Ref"], rotation=0)
    ax.set_ylabel("SI-SDR / dB")
    ax.set_title("Separation quality by SI-SDR")
    ax.grid(True, axis="y", alpha=0.45)
    ax.legend(frameon=False, loc="upper left")
    sisdr = ASSETS / "metric_sisdr.png"
    save_fig(sisdr)

    pivot_corr = metrics.pivot(index="method", columns="stem", values="corr").loc[order]
    fig, ax = plt.subplots(figsize=(8.6, 4.6))
    ax.plot(x, pivot_corr["male"], marker="o", linewidth=2.4, color=f"#{BLUE_DARK}", label="male")
    ax.plot(x, pivot_corr["female"], marker="o", linewidth=2.4, color=f"#{PINK}", label="female")
    ax.set_xticks(x)
    ax.set_xticklabels(["Butter", "Cheby-I", "Ellip", "FIR", "Ref", "Diff", "T-Ref"])
    ax.set_ylim(0.25, 1.02)
    ax.set_ylabel("Correlation")
    ax.set_title("Waveform similarity to the reference")
    ax.grid(True, axis="y", alpha=0.45)
    ax.legend(frameon=False, loc="lower right")
    corr = ASSETS / "metric_corr.png"
    save_fig(corr)

    return {"sisdr": sisdr, "corr": corr}


def create_audio_charts() -> dict[str, Path]:
    morandi_style()
    male, sr = read_audio(AUDIO / "male_speech_ref.wav")
    female, _ = read_audio(AUDIO / "female_speech_ref.wav")
    mix, _ = read_audio(AUDIO / "mix_speech.wav")
    tref_male, _ = read_audio(AUDIO / "tref_male.wav")
    tref_female, _ = read_audio(AUDIO / "tref_female.wav")
    butter_male, _ = read_audio(AUDIO / "butter_male.wav")
    butter_female, _ = read_audio(AUDIO / "butter_female.wav")

    t = np.arange(len(mix)) / sr
    fig, axes = plt.subplots(3, 1, figsize=(9.0, 5.3), sharex=True)
    for ax, y, title, color in [
        (axes[0], male, "Male reference", BLUE_DARK),
        (axes[1], female, "Female reference", PINK),
        (axes[2], mix, "Mixed speech", TEXT),
    ]:
        ax.plot(t, y, color=f"#{color}", linewidth=0.55)
        ax.set_title(title, loc="left", fontsize=11)
        ax.grid(True, alpha=0.35)
        ax.set_ylabel("Amp.")
    axes[2].set_xlabel("Time / s")
    mix_chart = ASSETS / "wave_mix_refs.png"
    save_fig(mix_chart)

    fig, axes = plt.subplots(2, 2, figsize=(9.2, 5.2), sharex=True)
    panels = [
        (axes[0, 0], butter_male, "Butterworth male", WARM_GRAY),
        (axes[0, 1], butter_female, "Butterworth female", TERR),
        (axes[1, 0], tref_male, "TrainableRef male", BLUE_DARK),
        (axes[1, 1], tref_female, "TrainableRef female", PINK),
    ]
    for ax, y, title, color in panels:
        ax.plot(t, y, color=f"#{color}", linewidth=0.55)
        ax.set_title(title, loc="left", fontsize=11)
        ax.grid(True, alpha=0.35)
        ax.set_xlabel("Time / s")
    sep_chart = ASSETS / "wave_method_compare.png"
    save_fig(sep_chart)

    return {"mix": mix_chart, "sep": sep_chart}


def create_process_charts() -> dict[str, Path]:
    morandi_style()
    fig, ax = plt.subplots(figsize=(9.5, 2.6))
    ax.axis("off")
    labels = ["female_speech.wav", "Resample", "Waveform + spectrum", "Lowest fs ≈ 8 kHz"]
    xs = [0.08, 0.34, 0.61, 0.86]
    for i, (x, label) in enumerate(zip(xs, labels)):
        ax.add_patch(plt.Rectangle((x - 0.10, 0.35), 0.20, 0.30, color=f"#{[BLUE, SAGE, PINK, TERR][i]}", alpha=0.9))
        ax.text(x, 0.50, label, ha="center", va="center", fontsize=11, color=f"#{TEXT}")
        if i < len(xs) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.12, 0.50), xytext=(x + 0.12, 0.50), arrowprops=dict(arrowstyle="->", lw=1.8, color=f"#{MUTED}"))
    p1 = ASSETS / "process_sampling.png"
    save_fig(p1)

    fig, ax = plt.subplots(figsize=(9.5, 3.0))
    ax.axis("off")
    labels = ["male + female", "mix_speech.wav", "Fixed filters", "Learnable masks", "Metrics"]
    xs = [0.08, 0.29, 0.50, 0.71, 0.91]
    for i, (x, label) in enumerate(zip(xs, labels)):
        ax.add_patch(plt.Circle((x, 0.52), 0.095, color=f"#{[BLUE, SAGE, WARM_GRAY, PINK, TERR][i]}", alpha=0.95))
        ax.text(x, 0.52, label, ha="center", va="center", fontsize=10, color=f"#{TEXT}")
        if i < len(xs) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.11, 0.52), xytext=(x + 0.11, 0.52), arrowprops=dict(arrowstyle="->", lw=1.7, color=f"#{MUTED}"))
    p2 = ASSETS / "process_separation.png"
    save_fig(p2)
    return {"sampling": p1, "separation": p2}


def create_all_assets() -> dict[str, Path]:
    ensure_assets()
    assets = {}
    assets.update(create_metric_charts())
    assets.update(create_audio_charts())
    assets.update(create_process_charts())
    return assets


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=TEXT, font="Microsoft YaHei", align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font
    return tb


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.72, 0.35, 9.5, 0.45, size=24, bold=True, color=TEXT, font="Cambria")
    if subtitle:
        add_text(slide, subtitle, 0.74, 0.83, 8.5, 0.28, size=10.5, color=MUTED)
    add_shape(slide, 0.42, 0.38, 0.12, 0.42, BLUE, radius=True)


def add_shape(slide, x, y, w, h, fill, line=None, radius=True, alpha=None):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def add_panel(slide, x, y, w, h, fill=PANEL):
    return add_shape(slide, x, y, w, h, fill=fill, line=LINE, radius=True)


def add_image(slide, path: Path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w) if w else None, height=Inches(h) if h else None)


def add_metric_card(slide, label, value, x, y, color):
    add_shape(slide, x, y, 2.25, 1.0, color, radius=True)
    add_text(slide, value, x + 0.18, y + 0.13, 1.9, 0.34, size=23, bold=True, color=TEXT, font="Cambria")
    add_text(slide, label, x + 0.18, y + 0.58, 1.9, 0.22, size=9.8, color=TEXT)


def add_bullets(slide, items, x, y, w, h, size=13.2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = rgb(TEXT)
        p.space_after = Pt(7)
    return tb


def build_deck(assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        add_bg(s)
        return s

    # 1
    s = slide()
    add_shape(s, 0, 0, W, H, BG, radius=False)
    add_shape(s, 0.65, 0.6, 4.7, 5.95, BLUE, radius=True)
    add_shape(s, 5.65, 0.95, 6.7, 4.7, PANEL, line=LINE, radius=True)
    add_text(s, "语音信号处理与男女声混合分离", 1.05, 1.10, 4.1, 1.65, size=30, bold=True, color=TEXT, font="Cambria")
    add_text(s, "Signal Analysis & Processing · Problem 4", 1.08, 2.95, 3.5, 0.35, size=12, color=TEXT)
    add_text(s, "采样频率分析 / 传统滤波器 / 可微滤波与参考说话人引导", 5.95, 1.30, 5.9, 0.45, size=17, bold=True)
    add_bullets(s, ["课程设计第四题汇报", "姓名 / 学号 / 班级", "数据：男声与女声语音样本"], 5.95, 2.2, 5.6, 1.6, size=15)
    add_metric_card(s, "最低不失真采样率", "8 kHz", 5.95, 4.35, SAGE)
    add_metric_card(s, "最佳 SI-SDR", "14.52 dB", 8.45, 4.35, PINK)

    # 2
    s = slide()
    add_title(s, "原题要求拆解", "Speech signal processing")
    cards = [
        ("1", "录制与播放", "录制一段 3 秒语音并播放，作为后续分析对象。"),
        ("2", "采样率分析", "使用不同采样频率，画出波形和频谱，找最低不失真采样率。"),
        ("3", "混合与分离", "分别录男/女语音，混合后设计滤波器分离并与原始语音比较。"),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = 0.8 + i * 4.1
        add_panel(s, x, 1.55, 3.5, 3.9)
        add_shape(s, x + 0.25, 1.85, 0.55, 0.55, [BLUE, SAGE, PINK][i], radius=True)
        add_text(s, num, x + 0.42, 1.96, 0.2, 0.2, size=16, bold=True, font="Cambria")
        add_text(s, head, x + 0.25, 2.62, 2.8, 0.35, size=18, bold=True)
        add_text(s, body, x + 0.25, 3.25, 2.9, 1.25, size=13, color=MUTED)

    # 3
    s = slide()
    add_title(s, "实验总体流程", "从采样分析到混合分离，再扩展到可学习方法")
    add_panel(s, 0.7, 1.25, 12.0, 2.25)
    add_image(s, assets["sampling"], 1.0, 1.68, w=11.4)
    add_panel(s, 0.7, 4.15, 12.0, 2.15)
    add_image(s, assets["separation"], 1.0, 4.58, w=11.4)

    # 4
    s = slide()
    add_title(s, "数据来源与混合构造", "所有比较均基于同一组有真值的男/女语音")
    add_panel(s, 0.7, 1.2, 6.0, 5.35)
    add_image(s, assets["mix"], 1.0, 1.55, w=5.4)
    add_panel(s, 7.05, 1.2, 5.55, 5.35)
    add_bullets(s, ["male_speech_ref.wav：男声参考，9.32 s", "female_speech_ref.wav：女声参考，补静音至 9.32 s", "mix_speech.wav = male + female", "因为有标准答案，可以计算 SI-SDR、相关系数和 SNR"], 7.45, 1.65, 4.65, 3.7, size=14)

    # 5
    s = slide()
    add_title(s, "第二问：不同采样频率实验设计", "固定语音内容，只改变采样频率")
    add_panel(s, 0.8, 1.35, 5.5, 4.8)
    add_text(s, "采样率设置", 1.15, 1.75, 2.2, 0.35, size=18, bold=True)
    add_bullets(s, ["16000 Hz：原始参考", "12000 Hz：中高质量采样", "8000 Hz：电话语音标准", "4000 Hz：高频明显损失", "2000 Hz：严重失真"], 1.15, 2.35, 4.4, 2.6, size=14)
    add_panel(s, 6.8, 1.35, 5.4, 4.8)
    add_text(s, "理论依据", 7.15, 1.75, 2.2, 0.35, size=18, bold=True)
    add_text(s, "语音主要有效频带约为 300–3400 Hz。\n根据奈奎斯特采样定理，采样频率需满足 fs ≥ 2fmax，因此最低不失真采样率约为 6.8 kHz，实际常取 8 kHz。", 7.15, 2.4, 4.4, 2.2, size=14, color=TEXT)
    add_metric_card(s, "工程取值", "8 kHz", 7.15, 5.0, BLUE)

    # 6
    s = slide()
    add_title(s, "第二问结果：波形与频谱", "采样率降低会压缩可表示频带")
    add_panel(s, 0.55, 1.1, 12.25, 5.95)
    add_image(s, IMAGES / "problem4_sampling_rate_comparison.png", 0.85, 1.38, w=11.65)

    # 7
    s = slide()
    add_title(s, "最低不失真采样率判断", "8 kHz 是语音通信中的实用下限")
    add_metric_card(s, "Nyquist 估算", "6.8 kHz", 0.95, 1.6, SAGE)
    add_metric_card(s, "实验判断", "8 kHz", 3.55, 1.6, BLUE)
    add_metric_card(s, "4 kHz 以下", "失真", 6.15, 1.6, PINK)
    add_panel(s, 0.95, 3.1, 11.3, 2.6)
    add_bullets(s, ["8 kHz 能保留主要语音信息，频谱覆盖到 4 kHz。", "4 kHz 采样的 Nyquist 频率只有 2 kHz，高频辅音、齿音和音色细节会丢失。", "2 kHz 采样只能覆盖 1 kHz 以下，语音可懂度与自然度都显著下降。"], 1.35, 3.55, 10.2, 1.8, size=15)

    # 8
    s = slide()
    add_title(s, "第三问：男女声混合与标准答案", "人工混合使客观评价成为可能")
    add_panel(s, 0.7, 1.25, 5.8, 4.9)
    add_text(s, "混合关系", 1.05, 1.65, 2.2, 0.35, size=18, bold=True)
    add_text(s, "mix_speech(t) = male_ref(t) + female_ref(t)", 1.05, 2.3, 4.75, 0.45, size=18, bold=True, color=BLUE_DARK, font="Cambria")
    add_bullets(s, ["两路参考语音长度统一为 9.32 秒", "女生后段补静音，保留男声完整内容", "所有方法输入同一个 mix_speech.wav"], 1.05, 3.15, 4.8, 1.8, size=14)
    add_panel(s, 6.85, 1.25, 5.6, 4.9)
    add_image(s, assets["mix"], 7.15, 1.65, w=5.0)

    # 9
    s = slide()
    add_title(s, "传统滤波器设计", "固定频率选择性滤波是原题核心方法")
    methods = [("Butterworth", "幅频响应平滑，通带无波纹"), ("Chebyshev-I", "过渡更陡，通带存在波纹"), ("Elliptic", "过渡最陡，通带/阻带均有波纹"), ("FIR-Hamming", "线性相位，阶数较高")]
    for i, (m, desc) in enumerate(methods):
        x = 0.8 + (i % 2) * 6.1
        y = 1.35 + (i // 2) * 2.25
        add_panel(s, x, y, 5.45, 1.65)
        add_shape(s, x + 0.25, y + 0.35, 0.55, 0.55, [BLUE, SAGE, PINK, TERR][i])
        add_text(s, m, x + 1.0, y + 0.32, 3.5, 0.3, size=17, bold=True)
        add_text(s, desc, x + 1.0, y + 0.85, 3.9, 0.28, size=12.5, color=MUTED)

    # 10
    s = slide()
    add_title(s, "滤波器幅频响应对比", "不同传统滤波器的频率选择性不同")
    add_panel(s, 0.65, 1.05, 12.1, 5.95)
    add_image(s, IMAGES / "problem4_speech_filter_response_comparison.png", 0.95, 1.35, w=11.5)

    # 11
    s = slide()
    add_title(s, "传统滤波分离结果与局限", "固定切频难以处理男女声频谱重叠")
    add_panel(s, 0.65, 1.05, 7.0, 5.8)
    add_image(s, assets["sep"], 0.95, 1.38, w=6.4)
    add_panel(s, 8.0, 1.05, 4.65, 5.8)
    add_bullets(s, ["低通能保留一部分男声主体，因此男声指标略高。", "女声不是只有高频，高通会丢失共振峰与主体能量。", "男声的高频泛音会进入女声输出，造成明显串音。"], 8.35, 1.55, 3.85, 2.4, size=13.5)
    add_metric_card(s, "最佳传统男声", "1.27 dB", 8.35, 4.35, SAGE)
    add_metric_card(s, "传统女声", "仍较差", 10.1, 4.35, PINK)

    # 12
    s = slide()
    add_title(s, "可微滤波模型", "频率先验不再是固定后处理，而是参与训练")
    add_panel(s, 0.75, 1.25, 5.6, 4.9)
    add_bullets(s, ["输入：mix_speech.wav", "STFT 得到时频表示", "神经网络预测 male/female mask", "加入可学习低通/高通先验", "iSTFT 还原两路语音"], 1.1, 1.7, 4.75, 3.1, size=14)
    add_panel(s, 6.75, 1.25, 5.65, 4.9)
    add_image(s, IMAGES / "problem4_speech_training_loss.png", 7.05, 1.58, w=5.05)

    # 13
    s = slide()
    add_title(s, "Reference-guided：参考说话人引导", "从“按频率分”升级到“按目标音色分”")
    add_panel(s, 0.7, 1.2, 5.75, 5.1)
    add_bullets(s, ["target_ref.wav 提供目标说话人的频谱模板", "mix_speech.wav 提供待分离混合语音", "模板相似度 + 时频 mask 决定每个频点归属", "TrainableRefGuided 进一步用真实源监督训练"], 1.05, 1.65, 4.8, 3.4, size=14)
    add_panel(s, 6.85, 1.2, 5.45, 5.1)
    add_image(s, IMAGES / "problem4_reference_training_loss.png", 7.15, 1.55, w=4.9)

    # 14
    s = slide()
    add_title(s, "客观指标对比", "SI-SDR 越高，分离结果越接近真实源")
    add_panel(s, 0.65, 1.05, 6.0, 5.75)
    add_image(s, assets["sisdr"], 0.95, 1.45, w=5.45)
    add_panel(s, 6.95, 1.05, 5.75, 5.75)
    add_image(s, assets["corr"], 7.25, 1.45, w=5.15)

    # 15
    s = slide()
    add_title(s, "最佳结果：TrainableRefGuided", "参考音色 + 监督训练带来最稳定的提升")
    add_metric_card(s, "男声 SI-SDR", "14.52 dB", 0.95, 1.35, BLUE)
    add_metric_card(s, "女声 SI-SDR", "14.50 dB", 3.55, 1.35, PINK)
    add_metric_card(s, "相关系数", "0.982+", 6.15, 1.35, SAGE)
    add_panel(s, 0.95, 3.0, 11.25, 2.75)
    add_bullets(s, ["固定滤波器只能按频率切分，难以解决男女声重叠。", "可微滤波模型能学习时频 mask，因此显著提升分离质量。", "TrainableRefGuided 利用目标说话人参考音色，在本实验中取得最佳指标。"], 1.35, 3.45, 10.0, 1.9, size=15)

    # 16
    s = slide()
    add_title(s, "结论", "传统 DSP 与可学习方法的互补")
    add_panel(s, 0.85, 1.25, 3.65, 4.65)
    add_text(s, "采样分析", 1.15, 1.65, 2.0, 0.35, size=18, bold=True)
    add_text(s, "8 kHz 基本满足语音不明显失真的最低要求。", 1.15, 2.35, 2.8, 1.0, size=14, color=MUTED)
    add_panel(s, 4.85, 1.25, 3.65, 4.65)
    add_text(s, "传统滤波", 5.15, 1.65, 2.0, 0.35, size=18, bold=True)
    add_text(s, "原理清晰、可解释，但固定切频无法解决频谱重叠。", 5.15, 2.35, 2.8, 1.1, size=14, color=MUTED)
    add_panel(s, 8.85, 1.25, 3.65, 4.65)
    add_text(s, "可学习方法", 9.15, 1.65, 2.0, 0.35, size=18, bold=True)
    add_text(s, "引入时频 mask 和参考音色后，分离指标明显提升。", 9.15, 2.35, 2.8, 1.1, size=14, color=MUTED)
    add_text(s, "Thank you", 5.25, 6.45, 2.8, 0.35, size=20, bold=True, color=BLUE_DARK, font="Cambria", align=PP_ALIGN.CENTER)

    prs.save(OUT)


def main() -> None:
    assets = create_all_assets()
    build_deck(assets)
    print(f"Saved PPTX: {OUT}")


if __name__ == "__main__":
    main()
